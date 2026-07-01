"""
generate_mobile_dcfc_slides.py
================================
Generates a 16:9 PowerPoint presentation for the Caltrans Mobile DCFC study.
Upload the output PPTX to Google Drive → open with Google Slides.

Slides:
  1. Title
  2. Deliverables Status
  3. Cost Model — XOS Hub MC02
  4. Cost Model — Kempower DGS
  5. Cost Model Gaps
  6. Fixed DCFC MILP Formulation
  7. Mobile Kempower — Same MILP, 3 Differences
  8. Mobile XOS — SOC Simulation vs MILP
  9. Formulation Comparison (all three)
  10. Next Steps & Timeline

Run:  python generate_mobile_dcfc_slides.py
"""
from __future__ import annotations
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ──────────────────────────────────────────────────────────────────────────────
# CONSTANTS
# ──────────────────────────────────────────────────────────────────────────────

OUT_PATH = Path(r"D:\Geotab_EV_Parameters\charger_sizing_test\mobile_dcfc_slides.pptx")

SW = Inches(13.33)   # slide width  (16:9)
SH = Inches(7.50)    # slide height

# Colors
DARK_BLUE    = RGBColor(0x1F, 0x38, 0x64)
MID_BLUE     = RGBColor(0x27, 0x5B, 0xA0)
LIGHT_BLUE   = RGBColor(0xD6, 0xE4, 0xF7)
PALE_BLUE    = RGBColor(0xEE, 0xF4, 0xFD)
ORANGE       = RGBColor(0xED, 0x7D, 0x31)
DARK_ORANGE  = RGBColor(0xC0, 0x50, 0x00)
GREEN        = RGBColor(0x37, 0x86, 0x48)
DARK_GREEN   = RGBColor(0x1E, 0x50, 0x27)
LT_GREEN     = RGBColor(0xE2, 0xEF, 0xDA)
LT_YELLOW    = RGBColor(0xFF, 0xEB, 0x9C)
LT_RED       = RGBColor(0xFF, 0xC7, 0xCE)
RED          = RGBColor(0xC0, 0x00, 0x00)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LT_GRAY      = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY     = RGBColor(0xD9, 0xD9, 0xD9)
DARK_GRAY    = RGBColor(0x40, 0x40, 0x40)
BLACK        = RGBColor(0x00, 0x00, 0x00)

XOS_BG       = RGBColor(0xEB, 0xF3, 0xFF)   # XOS column bg
KMP_BG       = RGBColor(0xFE, 0xF3, 0xE2)   # Kempower column bg
FIX_BG       = RGBColor(0xE8, 0xF5, 0xE9)   # Fixed DCFC column bg

FONT_BODY = "Calibri"
FONT_MONO = "Courier New"
TODAY     = date.today().strftime("%B %d, %Y")


# ──────────────────────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank


def rect(slide, left, top, width, height, fill, border=None):
    """Add a solid-colored rectangle (shape type 1 = rectangle)."""
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
    else:
        shp.line.fill.background()
    return shp


def txbox(slide, text, left, top, width, height,
          size=12, bold=False, color=DARK_GRAY, italic=False,
          align=PP_ALIGN.LEFT, font=FONT_BODY, wrap=True,
          line_space=None):
    """Add a plain text box."""
    box  = slide.shapes.add_textbox(left, top, width, height)
    tf   = box.text_frame
    tf.word_wrap = wrap
    p    = tf.paragraphs[0]
    p.alignment = align
    if line_space:
        from pptx.util import Pt as Pt2
        p.line_spacing = Pt2(line_space)
    run  = p.add_run()
    run.text         = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font
    return box


def header(slide, title, subtitle=None, bg=DARK_BLUE, ht=Inches(0.85)):
    rect(slide, Inches(0), Inches(0), SW, ht, bg)
    txbox(slide, title,
          Inches(0.25), Inches(0.08), Inches(12.8), ht,
          size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle,
              Inches(0.25), Inches(0.55), Inches(10), Inches(0.32),
              size=11, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)


def section_bar(slide, text, top, width=None, left=Inches(0.2),
                bg=MID_BLUE, fc=WHITE, fs=11):
    w = width or (SW - Inches(0.4))
    rect(slide, left, top, w, Inches(0.32), bg)
    txbox(slide, text,
          left + Inches(0.1), top + Inches(0.02), w - Inches(0.2), Inches(0.3),
          size=fs, bold=True, color=fc)


def footer(slide, text="Caltrans ZEV Infrastructure Study — IA 65A1281"):
    rect(slide, Inches(0), SH - Inches(0.28), SW, Inches(0.28), DARK_BLUE)
    txbox(slide, text,
          Inches(0.2), SH - Inches(0.26), Inches(10), Inches(0.25),
          size=9, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
    txbox(slide, TODAY,
          SW - Inches(1.8), SH - Inches(0.26), Inches(1.6), Inches(0.25),
          size=9, color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)


def status_badge(slide, label, left, top, w=Inches(1.3), h=Inches(0.28)):
    """Colored done/partial/not-done badge."""
    cfg = {
        "DONE":      (LT_GREEN, DARK_GREEN),
        "PARTIAL":   (LT_YELLOW, DARK_ORANGE),
        "NOT DONE":  (LT_RED, RED),
        "TBD":       (MED_GRAY, DARK_GRAY),
    }
    bg, fc = cfg.get(label.upper(), (MED_GRAY, DARK_GRAY))
    rect(slide, left, top, w, h, bg, border=fc)
    txbox(slide, label, left, top + Inches(0.03), w, h - Inches(0.04),
          size=9, bold=True, color=fc, align=PP_ALIGN.CENTER)


def simple_table(slide, headers, rows, left, top, width, height,
                 header_bg=DARK_BLUE, header_fc=WHITE,
                 row_bgs=None, col_widths=None, header_size=10, row_size=9):
    """Add a simple formatted table."""
    nrows = len(rows) + 1
    ncols = len(headers)

    if col_widths is None:
        cw = int(width / ncols)
        col_widths = [cw] * ncols

    tbl = slide.shapes.add_table(nrows, ncols, left, top, width, height).table

    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw

    # Header row
    for j, h_txt in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h_txt
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.text = h_txt
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = header_fc
        run.font.name = FONT_BODY

    # Data rows
    for i, row_data in enumerate(rows):
        bg = (row_bgs[i] if row_bgs and i < len(row_bgs)
              else (LT_GRAY if i % 2 == 0 else WHITE))
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.text = str(val)
            run.font.size = Pt(row_size)
            run.font.name = FONT_BODY
            run.font.color.rgb = BLACK

    return tbl


def mono_box(slide, text, left, top, width, height,
             bg=LT_GRAY, fs=9.5, fc=BLACK):
    """Monospace code/equation box with light background."""
    rect(slide, left, top, width, height, bg)
    txbox(slide, text, left + Inches(0.08), top + Inches(0.05),
          width - Inches(0.16), height - Inches(0.1),
          size=fs, color=fc, font=FONT_MONO, wrap=True)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ──────────────────────────────────────────────────────────────────────────────

def slide_title(prs):
    sl = blank_slide(prs)
    # Full-bleed dark blue top half
    rect(sl, Inches(0), Inches(0), SW, Inches(4.2), DARK_BLUE)
    # Gold accent stripe
    rect(sl, Inches(0), Inches(4.2), SW, Inches(0.08), ORANGE)
    # Bottom half light
    rect(sl, Inches(0), Inches(4.28), SW, SH - Inches(4.28), LT_GRAY)

    # Main title
    txbox(sl, "Mobile DC Fast Charging (MDCFC)\nCost Model & Optimization Formulation",
          Inches(0.5), Inches(0.8), Inches(12), Inches(2.5),
          size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    txbox(sl, "Caltrans ZEV Infrastructure — IA 65A1281\nXos Hub MC02  ·  Kempower DGS  ·  vs Fixed DCFC",
          Inches(0.5), Inches(3.3), Inches(12), Inches(0.9),
          size=16, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

    # Metadata box
    rect(sl, Inches(0.5), Inches(4.55), Inches(5.5), Inches(2.6), PALE_BLUE,
         border=MID_BLUE)
    info = (
        f"Date:      {TODAY}\n"
        "Sites:     4 Caltrans Maintenance Stations\n"
        "Scenarios: XOS-only  |  Kempower-only\n"
        "Status:    Cost model + formulation in progress\n"
        "Target:    Final configs & costs  ~July 1, 2026"
    )
    txbox(sl, info,
          Inches(0.7), Inches(4.7), Inches(5.1), Inches(2.3),
          size=12, color=DARK_BLUE, font=FONT_MONO)

    # Right side label
    txbox(sl, "Key deliverables:\n\n"
          "✓  Cost component breakdown\n"
          "✓  MILP formulation (Kempower)\n"
          "✓  SOC simulation formulation (XOS)\n"
          "○  Scenario analysis\n"
          "○  Sensitivity analysis\n"
          "○  Comparison vs fixed DCFC (report)",
          Inches(7), Inches(4.55), Inches(5.8), Inches(2.6),
          size=12, color=DARK_GRAY)

    footer(sl)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — DELIVERABLES STATUS
# ──────────────────────────────────────────────────────────────────────────────

def slide_deliverables(prs):
    sl = blank_slide(prs)
    header(sl, "Deliverables Status — Mobile DCFC (XOS & Kempower)",
           "Caltrans IA 65A1281  ·  As of " + TODAY)
    rect(sl, Inches(0), Inches(0.85), SW, SH - Inches(1.13), LT_GRAY)

    # Legend
    status_badge(sl, "DONE",     Inches(0.25), Inches(0.93))
    status_badge(sl, "PARTIAL",  Inches(1.6),  Inches(0.93))
    status_badge(sl, "NOT DONE", Inches(2.95), Inches(0.93))

    # Table data
    headers = ["Deliverable", "Component", "XOS", "Kempower", "Note"]
    col_w   = [Inches(3.2), Inches(3.2), Inches(1.3), Inches(1.3), Inches(3.8)]

    rows = [
        # Cost model breakdown
        ("Cost breakdown — purchase",       "Initial purchase",         "DONE",     "DONE",     "XOS: $245k (Caltrans quote). Kempower: DGS contract prices"),
        ("Cost breakdown — electrical",     "Electrical upgrade",        "DONE",     "DONE",     "XOS: $20–30k (480V). Kempower: $10–20k (240V). No trenching"),
        ("Cost breakdown — permit",         "Permit / inspection",       "DONE",     "DONE",     "~$2,000/unit assumed; CA commercial"),
        ("Cost breakdown — maintenance",    "Annual O&M",                "PARTIAL",  "DONE",     "XOS: $6k/yr assumed (no vendor quote yet). Kempower: $1,573/yr DGS"),
        ("Cost breakdown — warranty",       "Warranty cost",             "PARTIAL",  "NOT DONE", "XOS: FreeWire proxy ($30–35k/3yr). Kempower: no data yet"),
        ("Cost breakdown — energy",         "Depot utility / energy",    "NOT DONE", "NOT DONE", "Grid cost at depot not modeled for either scenario"),
        ("Cost breakdown — end-of-life",    "EOL / decommissioning",     "NOT DONE", "NOT DONE", "Research required — no cost data found yet"),
        # Scenario analysis
        ("Scenario analysis — usage",       "Usage pattern variations",  "PARTIAL",  "PARTIAL",  "Simulation/MILP runs per day; no compiled scenario report"),
        ("Scenario analysis — energy cost", "Energy cost variations",    "NOT DONE", "NOT DONE", "Parametric sweep on depot energy rate not built"),
        ("Scenario analysis — lifespan",    "Equipment lifespan sweep",  "NOT DONE", "NOT DONE", "Life fixed at 8yr/10yr; no sensitivity on this"),
        # Mobility value
        ("Mobility value analysis",         "Cost vs fixed DCFC (table)","PARTIAL",  "PARTIAL",  "Both tools exist; comparison report not compiled"),
        ("Mobility value analysis",         "Multi-site flex. deployment","NOT DONE","NOT DONE", "Single-site only; no cross-site utilization model"),
        ("Mobility value analysis",         "vs diesel equipment",        "NOT DONE","NOT DONE", "No diesel cost model"),
        # Sensitivity
        ("Sensitivity analysis",            "Charger cost ±X%",          "NOT DONE", "NOT DONE", "Meeting: flagged as 'later'"),
        ("Sensitivity analysis",            "Dwell time (1h→2h)",        "NOT DONE", "NOT DONE", "Meeting: flagged as 'later'"),
    ]

    status_colors = {
        "DONE":     LT_GREEN,
        "PARTIAL":  LT_YELLOW,
        "NOT DONE": LT_RED,
    }

    tbl = sl.shapes.add_table(len(rows) + 1, 5,
                              Inches(0.15), Inches(1.28),
                              SW - Inches(0.3), Inches(5.8)).table
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = cw

    hdrs = ["Deliverable", "Component", "XOS", "Kempower", "Note"]
    for j, h_txt in enumerate(hdrs):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h_txt
        run.font.size = Pt(9); run.font.bold = True
        run.font.color.rgb = WHITE; run.font.name = FONT_BODY

    for i, row in enumerate(rows):
        bg = LT_GRAY if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            if j in (2, 3):
                cell.fill.solid()
                cell.fill.fore_color.rgb = status_colors.get(val, MED_GRAY)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j in (2, 3) else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(8 if j == 4 else 9)
            run.font.bold = (j in (2, 3))
            run.font.color.rgb = (RED if val == "NOT DONE" else
                                  DARK_GREEN if val == "DONE" else
                                  DARK_ORANGE if val == "PARTIAL" else BLACK)
            run.font.name = FONT_BODY

    footer(sl)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — COST MODEL XOS
# ──────────────────────────────────────────────────────────────────────────────

def slide_cost_xos(prs):
    sl = blank_slide(prs)
    header(sl, "Cost Model — Xos Hub MC02 (Mobile DCFC)",
           "Per-unit lifecycle cost breakdown  |  1 port active per unit  |  280 kWh battery (TAI spec)")
    rect(sl, Inches(0), Inches(0.85), SW, SH - Inches(1.13), LT_GRAY)

    # LEFT: cost table
    section_bar(sl, "📋  Per-Unit Cost Components", Inches(0.93),
                width=Inches(6.5), bg=DARK_BLUE)

    headers = ["Component", "Amount (USD)", "Basis / Notes"]
    rows = [
        ("Purchase price",       "$245,437.50",         "Caltrans informal quote (per unit)"),
        ("Electrical upgrade",   "$20,000 – $30,000",   "480V 3-phase, no trenching, in-building only"),
        ("Permit / inspection",  "$2,000",               "CA commercial, assumed"),
        ("Annual maintenance",   "$6,000 / yr",          "Assumed — awaiting Xos vendor quote"),
        ("Warranty (proxy)",     "$3,000–3,500 / yr",    "FreeWire proxy ($30–35k / 3yr); update from Xos"),
        ("End-of-life",          "TBD",                  "Recycling / decommissioning — research needed"),
        ("Energy (depot)",       "Not modeled",          "Overnight grid charging cost not yet included"),
    ]
    row_bgs = [LT_GRAY, WHITE, LT_GRAY, WHITE, LT_GRAY, LT_RED, LT_RED]
    col_w2 = [Inches(1.6), Inches(1.7), Inches(3.1)]
    tbl = sl.shapes.add_table(len(rows)+1, 3,
                              Inches(0.2), Inches(1.28),
                              Inches(6.5), Inches(4.1)).table
    for i, cw in enumerate(col_w2):
        tbl.columns[i].width = cw
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = MID_BLUE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.size = Pt(9); run.font.bold = True
        run.font.color.rgb = WHITE; run.font.name = FONT_BODY
    for i, (comp, amt, note) in enumerate(rows):
        bg = row_bgs[i]
        for j, val in enumerate((comp, amt, note)):
            cell = tbl.cell(i+1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = val
            run.font.size = Pt(9); run.font.name = FONT_BODY
            run.font.bold = (j == 1)
            run.font.color.rgb = (RED if val in ("TBD", "Not modeled") else DARK_GRAY)

    # Lifecycle formula
    section_bar(sl, "💲  Lifecycle Cost Formula", Inches(5.52),
                width=Inches(6.5), bg=MID_BLUE)
    mono_box(sl,
             "C_lifecycle = purchase + electrical + permit\n"
             "            + maintenance × life_years\n"
             "            + warranty_total + EOL\n\n"
             "C_annual    = C_lifecycle / 10 years\n"
             "C_daily     = C_annual    / 365\n\n"
             "Low  estimate : $245,437 + $20,000 + $2,000 + $60,000\n"
             "                        + $30,000 + EOL  =  $357,437 + EOL\n"
             "High estimate : $245,437 + $30,000 + $2,000 + $60,000\n"
             "                        + $35,000 + EOL  =  $372,437 + EOL\n\n"
             "→  Annual (low / high): $35,744 / $37,244  per unit\n"
             "→  Daily  (low / high): $97.93  / $102.04  per unit",
             Inches(0.2), Inches(5.87), Inches(6.5), Inches(1.37),
             bg=PALE_BLUE, fs=9)

    # RIGHT: specs panel
    rect(sl, Inches(6.85), Inches(0.93), Inches(6.3), Inches(6.3), WHITE,
         border=MID_BLUE)
    section_bar(sl, "⚡  Xos Hub MC02 — Key Specifications",
                Inches(0.97), width=Inches(6.1), left=Inches(6.95), bg=MID_BLUE)

    specs = (
        "Battery capacity   :  282 kWh nominal  →  280 kWh (TAI spec)\n"
        "Usable energy      :  224 kWh  (SoC 20%–100%)\n"
        "Charge ports       :  4 × CCS1  (1 port modeled per unit)\n"
        "Port output        :  80 kW max (constant)\n"
        "Grid input rate    :  83 kW  (480V × 100A × √3)\n"
        "Grid connection    :  480V 3-phase  (no trenching)\n"
        "Battery chemistry  :  LFP (Lithium Iron Phosphate)\n"
        "Charge efficiency  :  η_c = 0.95\n"
        "Discharge effic.   :  η_d = 0.95\n"
        "Service life       :  10 years / 3,000 cycles @ 70% DoD\n"
    )
    txbox(sl, specs, Inches(7.0), Inches(1.35), Inches(6.0), Inches(2.4),
          size=10, font=FONT_MONO, color=DARK_BLUE)

    section_bar(sl, "🔄  Operational Model",
                Inches(3.83), width=Inches(6.1), left=Inches(6.95), bg=DARK_GRAY)
    ops = (
        "• XOS charges its battery from grid during off-peak / idle periods\n"
        "  (always recharges itself when not serving a vehicle)\n\n"
        "• During dispatch: battery discharges to vehicles (grid not involved)\n"
        "  → site grid peak = battery charging periods, NOT dispatch periods\n\n"
        "• Can also charge vehicles + accept grid simultaneously (up to 230 kW)\n"
        "  but primary model: off-peak charge → peak dispatch\n\n"
        "• Minimum SoC = 20%  maintained at all times (operational constraint)"
    )
    txbox(sl, ops, Inches(7.0), Inches(4.2), Inches(6.0), Inches(2.5),
          size=10, color=DARK_GRAY)

    footer(sl)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — COST MODEL KEMPOWER
# ──────────────────────────────────────────────────────────────────────────────

def slide_cost_kempower(prs):
    sl = blank_slide(prs)
    header(sl, "Cost Model — Kempower DGS (Mobile DCFC)",
           "CA DGS Contract 1-23-61-15A (National Car Charging LLC)  |  DC-only, grid-connected, no battery")
    rect(sl, Inches(0), Inches(0.85), SW, SH - Inches(1.13), LT_GRAY)

    # Three charger type columns
    types = [
        ("Kempower 50 kW",
         "Group 5  B-500/B-501",
         "$23,408",  "$855",   "$1,573/yr",  "8 yr",
         "$10–20k",  "$2,000", "$26,836", "$215"),
        ("Kempower 150 kW",
         "Group 6  S-600/S-601",
         "$62,154",  "$4,750", "$1,573/yr",  "8 yr",
         "$10–20k",  "$2,000", "$91,498", "$731"),
        ("Kempower 250 kW",
         "Group 7  S-700/S-701",
         "$101,946", "$5,225", "$1,573/yr",  "8 yr",
         "$10–20k",  "$2,000", "$133,755","$1,070"),
    ]

    col_bgs = [LT_GRAY, XOS_BG, KMP_BG]
    labels  = ["Purchase", "DGS Install", "Maint/yr", "Life",
               "Electrical", "Permit", "Lifecycle* (ex-EOL)", "Daily CapEx*"]

    for ci, (name, sub, pur, inst, maint, life, elec, perm, lc, daily) in enumerate(types):
        x = Inches(0.2) + ci * Inches(4.37)
        w = Inches(4.3)

        rect(sl, x, Inches(0.93), w, Inches(0.35), DARK_BLUE)
        txbox(sl, name, x + Inches(0.05), Inches(0.95), w - Inches(0.1), Inches(0.32),
              size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        rect(sl, x, Inches(1.28), w, Inches(0.22), MID_BLUE)
        txbox(sl, sub, x + Inches(0.05), Inches(1.29), w - Inches(0.1), Inches(0.21),
              size=9, color=WHITE, align=PP_ALIGN.CENTER)

        vals = [pur, inst, maint, life, elec, perm, lc, daily]
        for ri, (lbl, val) in enumerate(zip(labels, vals)):
            row_top = Inches(1.5) + ri * Inches(0.43)
            bg = LT_GRAY if ri % 2 == 0 else WHITE
            if lbl in ("Lifecycle* (ex-EOL)", "Daily CapEx*"):
                bg = LIGHT_BLUE
            rect(sl, x, row_top, w, Inches(0.42), bg)
            txbox(sl, lbl, x + Inches(0.05), row_top + Inches(0.04),
                  Inches(1.6), Inches(0.38), size=9, color=DARK_GRAY)
            txbox(sl, val, x + Inches(1.65), row_top + Inches(0.04),
                  w - Inches(1.7), Inches(0.38), size=10, bold=True,
                  color=DARK_BLUE, align=PP_ALIGN.RIGHT)

    # Bottom notes
    section_bar(sl, "📌  Notes", Inches(4.97), bg=DARK_GRAY)
    notes = (
        "* Lifecycle = purchase + DGS install + electrical ($15k mid) + permit + maintenance × 8yr    "
        "   Daily CapEx = [(purchase+install)/(life×12) + maint/12] / 30.42\n"
        "• End-of-life cost: TBD (research required — charger recycling / disposal)\n"
        "• Depot energy (grid) cost not modeled — add utility tariff at deployment site\n"
        "• No warranty data for Kempower yet — obtain from vendor\n"
        "• DGS allows 1% discount on orders >$100k, 2% on orders >$500k (pre-discount prices shown)"
    )
    txbox(sl, notes, Inches(0.2), Inches(5.2), SW - Inches(0.4), Inches(1.75),
          size=9, color=DARK_GRAY)

    footer(sl)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — FIXED DCFC MILP FORMULATION
# ──────────────────────────────────────────────────────────────────────────────

def slide_fixed_milp(prs):
    sl = blank_slide(prs)
    header(sl, "Fixed DCFC — MILP Formulation  (reference baseline)",
           "exact_northgate_charger_sizing_milp.py  |  Gurobi / HiGHS  |  Northgate site, one day at a time")
    rect(sl, Inches(0), Inches(0.85), SW, SH - Inches(1.13), LT_GRAY)

    # LEFT COLUMN: Sets, Vars
    section_bar(sl, "SETS & PARAMETERS", Inches(0.93), width=Inches(6.3), bg=DARK_BLUE)
    mono_box(sl,
             "V   vehicle charging events (one per depot visit)\n"
             "T   discrete 15-min time steps  (Δt = 5–15 min)\n"
             "C   charger types: {L2·19.2 kW, DC·50 kW, DC·150 kW, DC·350 kW}\n\n"
             "Ev  energy required by vehicle v        [kWh]\n"
             "Peff[v,c] = min(Pc, Pdc_max_v) or min(Pc, Pac_max_v)  [kW]\n"
             "η   = 0.90  (charging efficiency, grid→battery)\n"
             "Cdaily,c = [(purchase+install)/(life×12) + maint/12] / 30.42",
             Inches(0.2), Inches(1.28), Inches(6.3), Inches(1.6), fs=9.5)

    section_bar(sl, "DECISION VARIABLES", Inches(2.95), width=Inches(6.3), bg=DARK_BLUE)
    mono_box(sl,
             "Nc         ∈ ℤ+    number of chargers of type c installed\n"
             "u[v,t,c]   ∈ {0,1} 1 if vehicle v charges on type c at step t\n"
             "z[v,c]     ∈ {0,1} 1 if vehicle v is assigned to charger type c\n"
             "P_total[t] ∈ ℝ+    total site charging power at step t  [kW]\n"
             "P_max      ∈ ℝ+    peak site demand over full day  [kW]\n"
             "P_peak_win ∈ ℝ+    peak demand during 4–9 PM  [kW]",
             Inches(0.2), Inches(3.3), Inches(6.3), Inches(1.45), fs=9.5)

    section_bar(sl, "OBJECTIVE  —  minimize total daily cost", Inches(4.82),
                width=Inches(6.3), bg=DARK_BLUE)
    mono_box(sl,
             "min  Σc Nc · Cdaily,c                        ← charger CapEx (amortized)\n"
             "   + P_max · 6.45                            ← SMUD global demand charge\n"
             "   + P_peak_win · 9.96                       ← SMUD peak-window demand\n"
             "   + (λ/T) · Σt [max(0, P(t)−P̄)]²          ← above-mean power penalty\n"
             "   + Σt P(t) · cenergy(t) · Δt               ← SMUD TOD energy cost",
             Inches(0.2), Inches(5.17), Inches(6.3), Inches(1.45), fs=9.5)

    # RIGHT COLUMN: Constraints
    section_bar(sl, "KEY CONSTRAINTS", Inches(0.93), width=Inches(6.5),
                left=Inches(6.65), bg=DARK_GREEN)
    mono_box(sl,
             "(A) Energy lower bound (must fully charge each vehicle):\n"
             "     Σ{t,c} u[v,t,c] · Peff[v,c] · Δt · η  ≥  Ev     ∀v\n\n"
             "(B) Energy upper bound (no overcharge beyond battery room):\n"
             "     Σ{t,c} u[v,t,c] · Peff[v,c] · Δt · η  ≤  Emax_v  ∀v\n\n"
             "(C) Single-plug rule (one charger type active per vehicle per step):\n"
             "     Σc u[v,t,c]  ≤  1                                  ∀v,t\n\n"
             "(D) Charger capacity (active vehicles ≤ installed chargers):\n"
             "     Σv u[v,t,c]  ≤  Nc                                 ∀t,c\n\n"
             "(E) Charger exclusivity (vehicle uses one type for whole session):\n"
             "     Σc z[v,c]  ≤  1                                    ∀v\n"
             "     u[v,t,c]  ≤  z[v,c]                               ∀v,t,c\n\n"
             "(F) Contiguous charging (no replug within a session):\n"
             "     At most one 0→1 transition in {u[v,t,c] : t ∈ T}  ∀v,c\n"
             "     Enforced via binary session-start variable p[v,c,i]",
             Inches(6.65), Inches(1.28), Inches(6.5), Inches(3.9), fs=9.5)

    section_bar(sl, "SOLVER & OUTPUTS", Inches(5.25),
                width=Inches(6.5), left=Inches(6.65), bg=DARK_GRAY)
    txbox(sl,
          "Solver: Gurobi (primary)  ·  Pyomo + HiGHS (fallback)\n"
          "Output: N_c per type, charging schedule u[v,t,c], power profile P(t),\n"
          "        peak power P_max, cost breakdown by component\n"
          "Site:   Northgate (run for each day; top-30 configs → 90-95% coverage)",
          Inches(6.7), Inches(5.59), Inches(6.3), Inches(0.95),
          size=10, color=DARK_GRAY)

    footer(sl)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — MOBILE KEMPOWER: SAME MILP, 3 DIFFERENCES
# ──────────────────────────────────────────────────────────────────────────────

def slide_kempower_milp(prs):
    sl = blank_slide(prs)
    header(sl, "Mobile Kempower — Same MILP,  3 Key Differences vs Fixed DCFC",
           "kempower_milp_sizing.py  |  Gurobi / HiGHS  |  Wrapper around existing MILP")
    rect(sl, Inches(0), Inches(0.85), SW, SH - Inches(1.13), LT_GRAY)

    # "Same" banner
    rect(sl, Inches(0.2), Inches(0.93), Inches(8.7), Inches(0.35), LT_GREEN, border=GREEN)
    txbox(sl, "✓  IDENTICAL  to Fixed DCFC MILP:  "
          "all variables, all constraints (A–F), solver, contiguous-charging, "
          "energy bounds, single-plug, charger exclusivity",
          Inches(0.3), Inches(0.95), Inches(8.5), Inches(0.32),
          size=10, bold=False, color=DARK_GREEN)

    # THREE DIFFERENCE BOXES
    diffs = [
        ("① Charger Set  C",
         DARK_BLUE, PALE_BLUE,
         "Fixed DCFC:\n  {L2·19.2 kW, DC·50 kW, DC·150 kW, DC·350 kW}\n\n"
         "Kempower Mobile:\n  {Kempower·50 kW, Kempower·150 kW, Kempower·250 kW}\n\n"
         "→ No L2 chargers\n"
         "→ Vehicles with max_dc_charge_kw = 0 are excluded\n"
         "   (AC-only vehicles cannot use Kempower)"),
        ("② Cost Structure  Cdaily,c",
         DARK_ORANGE, KMP_BG,
         "Fixed DCFC  (Caltrans Q3 FY25/26):\n"
         "  DC·150: purchase $67.5k + install $77.5k  → higher install (trenching)\n\n"
         "Kempower DGS  (no trenching):\n"
         "  Kempower·150: purchase $62.2k + DGS install $4.75k\n"
         "  + electrical upgrade $10–20k + permit $2k\n"
         "  maint $1,573/yr (ChargerHelp! DGS rate, 8-yr life)\n\n"
         "→ No trenching → lower total install vs fixed DCFC\n"
         "→ Different daily CapEx Cdaily,c in objective"),
        ("③ Demand Charges",
         DARK_GREEN, LT_GREEN,
         "Fixed DCFC  (SMUD C&I 21–299kW):\n"
         "  Global demand:  P_max × $6.45/kW\n"
         "  Peak-window:    P_peak_win × $9.96/kW  (4–9 PM)\n\n"
         "Kempower Mobile:\n"
         "  Grid draw = at dispatch time (no battery buffer)\n"
         "  Demand charges depend on depot utility tariff\n"
         "  (Not necessarily SMUD — update for deployment site)\n\n"
         "→ Demand charge coefficients must match depot utility"),
    ]

    for i, (title, hdr_c, bg_c, body) in enumerate(diffs):
        x = Inches(0.2) + i * Inches(4.38)
        w = Inches(4.3)
        rect(sl, x, Inches(1.34), w, Inches(0.32), hdr_c)
        txbox(sl, title, x + Inches(0.08), Inches(1.36), w - Inches(0.16), Inches(0.3),
              size=11, bold=True, color=WHITE)
        mono_box(sl, body, x, Inches(1.66), w, Inches(3.5), bg=bg_c, fs=9)

    # Formulation mapping
    section_bar(sl, "MILP variables & constraints — unchanged (copy-paste from Fixed DCFC formulation)",
                Inches(5.22), bg=DARK_GRAY)
    txbox(sl,
          "Nc ∈ ℤ+  |  u[v,t,c] ∈ {0,1}  |  z[v,c] ∈ {0,1}  |  "
          "P_total[t]  |  P_max  |  P_peak_win\n"
          "Constraints A (energy lb)  ·  B (energy ub)  ·  C (single-plug)  ·  "
          "D (capacity)  ·  E (exclusivity)  ·  F (contiguous)",
          Inches(0.25), Inches(5.58), SW - Inches(0.5), Inches(0.6),
          size=10, color=DARK_GRAY, font=FONT_MONO)

    section_bar(sl, "IMPLEMENTATION NOTE", Inches(6.24), bg=MID_BLUE)
    txbox(sl,
          "kempower_milp_sizing.py patches 3 module globals then calls milp.main(charger_specs_override=kempower_specs):\n"
          "  milp.CHARGER_UPPER_BOUNDS = {Kempower_50kW: 20, Kempower_150kW: 15, Kempower_250kW: 10}\n"
          "  milp.OUTPUT_DIR = kempower_milp_outputs/\n"
          "  milp.INPUT_PATH_PRIMARY = <events_csv>",
          Inches(0.25), Inches(6.58), SW - Inches(0.5), Inches(0.62),
          size=9.5, color=DARK_BLUE, font=FONT_MONO)

    footer(sl)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — MOBILE XOS: SOC SIMULATION vs MILP
# ──────────────────────────────────────────────────────────────────────────────

def slide_xos_simulation(prs):
    sl = blank_slide(prs)
    header(sl, "Mobile XOS — SOC Simulation  (NOT an optimization)",
           "xos_hub_soc_simulation.py  |  Greedy time-series dispatch  |  Meeting decision: no MILP for XOS")
    rect(sl, Inches(0), Inches(0.85), SW, SH - Inches(1.13), LT_GRAY)

    # Why not MILP banner
    rect(sl, Inches(0.2), Inches(0.93), Inches(12.9), Inches(0.37), LT_YELLOW, border=ORANGE)
    txbox(sl, "⚠  Meeting decision (Jun 16 2026): No optimization / MILP for XOS.  "
          "Use a time-series SOC simulation — defensible, common-sense, no solver required.",
          Inches(0.3), Inches(0.96), Inches(12.6), Inches(0.33),
          size=10, bold=True, color=DARK_ORANGE)

    # LEFT: formulation
    section_bar(sl, "STATE VARIABLE & PARAMETERS", Inches(1.37), width=Inches(6.55), bg=DARK_BLUE)
    mono_box(sl,
             "K          number of XOS units deployed\n"
             "B  = 280 kWh      battery capacity (TAI spec)\n"
             "SOCmin = 20%      operational floor (XOS manual)\n"
             "Pgrid  = 83 kW    grid-to-battery input (480V·100A·√3)\n"
             "Pport  = 80 kW    battery-to-vehicle output (1 port per unit)\n"
             "ηc = ηd = 0.95   charge / discharge efficiency\n"
             "Δt = 0.25 h      time step (15 min)\n\n"
             "SOCk[t]  ∈ [20%, 100%]   — battery SoC of unit k at step t",
             Inches(0.2), Inches(1.72), Inches(6.55), Inches(1.9), fs=9.5)

    section_bar(sl, "STATE TRANSITION AT EACH STEP t", Inches(3.69), width=Inches(6.55), bg=DARK_BLUE)
    mono_box(sl,
             "Case A — unit k serves vehicle v:\n"
             "  energy to vehicle  = Peff,v · Δt · ηd         [kWh]\n"
             "  SOCk[t+1] = SOCk[t]  −  Peff,v · Δt / (ηd · B)\n"
             "  SOCk[t+1] = max(SOCk[t+1],  SOCmin)   [hard floor]\n\n"
             "Case B — unit k is idle (plugged in):\n"
             "  SOCk[t+1] = SOCk[t]  +  Pgrid · ηc · Δt / B\n"
             "  SOCk[t+1] = min(SOCk[t+1],  SOCmax)   [hard ceiling]\n\n"
             "  Peff,v = min(Pport, Pdc_max_v)   [kW at vehicle connector]",
             Inches(0.2), Inches(4.04), Inches(6.55), Inches(1.9), fs=9.5)

    section_bar(sl, "COVERAGE CRITERION", Inches(6.01), width=Inches(6.55), bg=DARK_GRAY)
    mono_box(sl,
             "Vehicle v is served  ⟺  remaining_v ≤ ε = 0.10 kWh  at departure",
             Inches(0.2), Inches(6.36), Inches(6.55), Inches(0.4), fs=9.5)

    # RIGHT: Greedy dispatch + sizing
    section_bar(sl, "GREEDY DISPATCH ALGORITHM  (each step t)", Inches(1.37),
                width=Inches(6.3), left=Inches(6.85), bg=DARK_GREEN)
    mono_box(sl,
             "1. Active vehicles:  A(t) = {v : av < t+Δt  AND  dv > t\n"
             "                             AND  rem_v > ε  AND  Peff,v > 0}\n\n"
             "2. Sort A(t) by urgency (descending):\n"
             "   urg_v = rem_v / max(dv − t,  Δt)      [kW — required rate]\n\n"
             "3. Available units:  U(t) = {k : not yet assigned this step}\n\n"
             "4. For each v ∈ A(t) by urgency:\n"
             "   (a) If U(t) empty → stop\n"
             "   (b) k* = argmax_{k ∈ U(t)} SOCk[t]     [highest SOC first]\n"
             "   (c) usable = (SOCk* − SOCmin) · B · ηd  [kWh to vehicle]\n"
             "   (d) If usable < ε → skip this unit\n"
             "   (e) e_del = min(Peff,v · Δt · ηd,  rem_v,  usable)\n"
             "   (f) Apply Case A for (k*, v);  U(t) ← U(t) \\ {k*}\n\n"
             "5. Apply Case B (grid charge) for all k ∉ serving(t)",
             Inches(6.85), Inches(1.72), Inches(6.3), Inches(3.3), fs=9.5)

    section_bar(sl, "ADD-ONE-UNTIL-COVERED  —  Sizing Rule",
                Inches(5.09), width=Inches(6.3), left=Inches(6.85), bg=MID_BLUE)
    mono_box(sl,
             "For K = 1, 2, 3, … (up to Kmax = 10):\n"
             "  1. Initialize:  SOCk[0] = 100%  for all k = 0..K-1\n"
             "  2. Run greedy dispatch simulation for full day\n"
             "  3. If ∀v: rem_v ≤ ε  →  return K  (minimum sufficient)\n\n"
             "Output: K units, SOC time-series, energy delivered per vehicle,\n"
             "        peak dispatch power, peak grid draw for battery charging",
             Inches(6.85), Inches(5.44), Inches(6.3), Inches(1.35), fs=9.5)

    footer(sl)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — SIDE-BY-SIDE FORMULATION COMPARISON
# ──────────────────────────────────────────────────────────────────────────────

def slide_comparison(prs):
    sl = blank_slide(prs)
    header(sl, "Formulation Comparison — Fixed DCFC vs Kempower vs XOS",
           "What changes and what stays the same across the three scenarios")
    rect(sl, Inches(0), Inches(0.85), SW, SH - Inches(1.13), LT_GRAY)

    # Column headers
    cols = [("Fixed DCFC", DARK_GREEN, FIX_BG),
            ("Mobile Kempower", DARK_ORANGE, KMP_BG),
            ("Mobile XOS", MID_BLUE, XOS_BG)]
    for i, (label, hc, _) in enumerate(cols):
        x = Inches(2.55) + i * Inches(3.55)
        rect(sl, x, Inches(0.93), Inches(3.48), Inches(0.35), hc)
        txbox(sl, label, x + Inches(0.05), Inches(0.95), Inches(3.38), Inches(0.32),
              size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Row labels column
    row_labels = [
        "Approach",
        "Solver",
        "Decision vars",
        "Charger types",
        "Battery / storage",
        "Energy constraints",
        "Demand charges",
        "Cost structure",
        "Sizing output",
        "Key assumption",
    ]
    rows_data = [
        ("MILP optimization",           "MILP optimization",         "Greedy simulation"),
        ("Gurobi / HiGHS",              "Gurobi / HiGHS",            "None (no solver)"),
        ("Nc, u[v,t,c], z[v,c],\nP_total, P_max, P_peak",
         "Nc, u[v,t,c], z[v,c],\nP_total, P_max, P_peak",
         "SOCk[t]\n(continuous, per unit)"),
        ("L2·19.2 / DC·50/150/350 kW",  "DC·50/150/250 kW (Kempower)\nNo L2",
         "80 kW per port\n(Xos Hub MC02)"),
        ("No  (grid-direct)",           "No  (grid-direct)",         "Yes  — 280 kWh LFP\n(1 port per unit)"),
        ("A: Ev ≤ Σ energy ≤ Emax\nC: single-plug\nD: capacity\nE: exclusivity\nF: contiguous",
         "Same A–F\n(unchanged)",
         "rem_v tracked per step\n(no explicit constraint —\nchecked at departure)"),
        ("SMUD C&I rates\n($6.45+$9.96/kW)",
         "Depot utility tariff\n(update for site)",
         "No demand charges\n(XOS charges off-peak;\ngrid draw = idle periods)"),
        ("Caltrans Q3 FY25/26\n(trenching included)",
         "DGS contract\n+ electrical $10–20k\n+ permit $2k (no trenching)",
         "Purchase $245k\n+ electrical $20–30k\n+ permit + maint + warranty"),
        ("N_c per type\n(minimum-cost mix)",
         "N_c per type\n(minimum-cost mix)",
         "K units\n(add-one-until-covered)"),
        ("No battery → always draws\nfrom grid at dispatch time",
         "No battery → always draws\nfrom grid at dispatch time",
         "Battery → grid draw ONLY\nduring idle (off-peak charge)"),
    ]

    row_h = Inches(0.54)
    for ri, (lbl, row_vals) in enumerate(zip(row_labels, rows_data)):
        y = Inches(1.34) + ri * row_h
        bg = LT_GRAY if ri % 2 == 0 else WHITE

        # Label cell
        rect(sl, Inches(0.15), y, Inches(2.35), row_h, DARK_BLUE if ri == 0 else MED_GRAY)
        txbox(sl, lbl, Inches(0.2), y + Inches(0.04), Inches(2.25), row_h - Inches(0.08),
              size=9, bold=True,
              color=(WHITE if ri == 0 else DARK_GRAY))

        # Data cells
        same_01 = (row_vals[0] == row_vals[1]) or ("Same" in str(row_vals[1]))
        same_12 = (row_vals[1] == row_vals[2])

        for ci, val in enumerate(row_vals):
            x = Inches(2.55) + ci * Inches(3.55)
            cell_bg = FIX_BG if ci == 0 else KMP_BG if ci == 1 else XOS_BG
            if ri % 2 == 0:
                pass  # keep scenario color
            else:
                cell_bg = bg

            # Highlight "same" cells
            if ci == 1 and same_01:
                cell_bg = LT_GREEN
            if ci == 2 and same_12:
                cell_bg = LT_GREEN

            rect(sl, x, y, Inches(3.48), row_h, cell_bg)
            txbox(sl, val, x + Inches(0.05), y + Inches(0.03),
                  Inches(3.38), row_h - Inches(0.05),
                  size=8.5, color=(DARK_GREEN if "Same" in val else DARK_GRAY),
                  font=FONT_BODY)

    # Legend
    rect(sl, SW - Inches(2.4), Inches(0.93), Inches(2.2), Inches(0.35), LT_GRAY, border=MED_GRAY)
    txbox(sl, "🟩 = same as Fixed DCFC",
          SW - Inches(2.35), Inches(0.97), Inches(2.1), Inches(0.28),
          size=8.5, color=DARK_GREEN)

    footer(sl)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — COST COMPARISON: MOBILE vs FIXED (GAP ANALYSIS)
# ──────────────────────────────────────────────────────────────────────────────

def slide_cost_comparison(prs):
    sl = blank_slide(prs)
    header(sl, "Cost Comparison — Mobile DCFC vs Fixed DCFC  (framework ready, report not compiled)",
           "What we can already compute vs. what requires further work")
    rect(sl, Inches(0), Inches(0.85), SW, SH - Inches(1.13), LT_GRAY)

    # LEFT: What's available now
    section_bar(sl, "✓  AVAILABLE NOW — Can be computed immediately", Inches(0.93),
                width=Inches(6.2), bg=DARK_GREEN)
    avail = (
        "Per-unit cost (partially complete):\n"
        "  • XOS:      $357k–$372k lifecycle  (excl. EOL + depot energy)\n"
        "  • Kempower: $100–135k lifecycle    (excl. EOL + depot energy)\n"
        "  • Fixed DC·150: ~$145k installed   (Caltrans Quarterly Q3)\n\n"
        "Sizing output (run the scripts):\n"
        "  • XOS:     K units from SOC simulation  →  K × per-unit cost\n"
        "  • Kempower: N_c units from MILP          →  Σ N_c × per-unit cost\n"
        "  • Fixed:    N_c units from MILP           →  Σ N_c × per-unit cost\n\n"
        "Infrastructure savings (no trenching):\n"
        "  Fixed DC·150 install: ~$77.5k  (trenching included)\n"
        "  Kempower install:     ~$4.75k + $15k elec = ~$20k  (no trench)\n"
        "  → Saving per 150 kW unit:  ~$57k in install cost\n\n"
        "Grid interaction:\n"
        "  Fixed DCFC:  draws from grid at every dispatch step\n"
        "  XOS:         grid draw only during off-peak charging → load-shift value"
    )
    txbox(sl, avail, Inches(0.2), Inches(1.28), Inches(6.2), Inches(5.0),
          size=10, color=DARK_GRAY)

    # RIGHT: What's missing
    section_bar(sl, "✗  MISSING — Required for complete deliverable", Inches(0.93),
                width=Inches(6.55), left=Inches(6.6), bg=RED)
    missing = [
        ("End-of-life cost",
         "Battery recycling / charger disposal cost for XOS and Kempower.\n"
         "Research queries: 'EV charger end-of-life cost', 'LFP battery recycling per kWh'"),
        ("Depot energy cost",
         "Utility tariff at each depot for grid charging.\n"
         "XOS charges overnight (off-peak); Kempower draws at dispatch.\n"
         "Need depot utility rate to compute annual energy cost."),
        ("Warranty (Kempower)",
         "No warranty data for Kempower DGS units yet.\n"
         "XOS: use FreeWire proxy ($30–35k / 3yr) until Xos rep responds."),
        ("Mobility value — quantified",
         "Fleet utilization across multiple sites, reduced downtime vs fixed,\n"
         "flexible redeployment when fleet size changes.  Qualitative only now."),
        ("vs Diesel comparison",
         "Diesel generator cost model not built.  No fuel, maintenance, or\n"
         "emissions data for mobile diesel charging alternative."),
        ("Scenario + sensitivity reports",
         "Parametric sweeps: charger cost ±50%, dwell time 1h→2h,\n"
         "energy cost variation, equipment lifespan 6–12yr.\n"
         "Meeting: 'later' — but required for full deliverable."),
    ]
    y = Inches(1.3)
    for i, (title, desc) in enumerate(missing):
        bg = LT_RED if i % 2 == 0 else WHITE
        rect(sl, Inches(6.6), y, Inches(6.55), Inches(0.8), bg)
        txbox(sl, f"✗  {title}", Inches(6.65), y + Inches(0.02),
              Inches(6.4), Inches(0.25), size=10, bold=True, color=RED)
        txbox(sl, desc, Inches(6.65), y + Inches(0.27),
              Inches(6.4), Inches(0.52), size=9, color=DARK_GRAY)
        y += Inches(0.82)

    footer(sl)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — NEXT STEPS & TIMELINE
# ──────────────────────────────────────────────────────────────────────────────

def slide_next_steps(prs):
    sl = blank_slide(prs)
    header(sl, "Next Steps & Timeline",
           "Target: finalize permanent + mobile charger costs and configs by ~July 1, 2026")
    rect(sl, Inches(0), Inches(0.85), SW, SH - Inches(1.13), LT_GRAY)

    steps = [
        # (priority, deadline, who, task, status_color)
        ("🔴 HIGH", "This week",   "You",        "Run XOS simulation + Kempower MILP on Northgate events → get unit counts",   GREEN),
        ("🔴 HIGH", "~Jun 23",     "You",        "Write mobile DCFC formulation section for quarterly report (use docstring from xos_hub_soc_simulation.py)", ORANGE),
        ("🔴 HIGH", "~Jun 23",     "You",        "Get EOL cost: search 'EV charger decommissioning cost', 'LFP battery recycling per kWh'", ORANGE),
        ("🔴 HIGH", "~Jun 23",     "Vendor",     "XOS rep: get actual warranty / maintenance contract pricing",                 ORANGE),
        ("🟠 MED",  "~Jun 27",     "You",        "Update permanent charger costs: 150/350 kW averages, 50 kW median, L2 80A classification", ORANGE),
        ("🟠 MED",  "~Jun 27",     "You",        "Add depot energy cost to XOS and Kempower models (get depot utility tariff)", ORANGE),
        ("🟠 MED",  "~Jun 30",     "You",        "Run fixed DCFC MILP for all 4 sites, every day → top-30 configs → 90-95% coverage", ORANGE),
        ("🟡 LOW",  "~Jul 7",      "You",        "Compile cost comparison table: XOS vs Kempower vs Fixed DCFC vs Diesel",     RED),
        ("🟡 LOW",  "~Jul 7",      "You",        "Scenario analysis: different usage days, energy cost ±20%",                  RED),
        ("🟡 LOW",  "Later",       "You",        "Sensitivity: dwell time 1h→2h, charger cost ±50%, lifespan 6–12yr",          RED),
    ]

    headers = ["Priority", "By", "Owner", "Task", "Status"]
    col_w2 = [Inches(0.85), Inches(0.85), Inches(0.7), Inches(9.3), Inches(0.8)]
    tbl = sl.shapes.add_table(len(steps)+1, 5,
                              Inches(0.15), Inches(0.93),
                              SW - Inches(0.3), Inches(5.9)).table
    for i, cw in enumerate(col_w2):
        tbl.columns[i].width = cw

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.size = Pt(10); run.font.bold = True
        run.font.color.rgb = WHITE; run.font.name = FONT_BODY

    for i, (pri, by, who, task, sc) in enumerate(steps):
        bg = LT_GRAY if i % 2 == 0 else WHITE
        vals = [pri, by, who, task, ""]
        for j, val in enumerate(vals):
            cell = tbl.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (sc if j == 4 else bg)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j != 3 else PP_ALIGN.LEFT
            run = p.add_run(); run.text = val
            run.font.size = Pt(9 if j != 3 else 9.5)
            run.font.name = FONT_BODY
            run.font.bold = (j == 0)
            run.font.color.rgb = (WHITE if j == 4 else
                                  RED if "HIGH" in val else
                                  DARK_ORANGE if "MED" in val else
                                  DARK_GRAY)

    footer(sl)


# ──────────────────────────────────────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("Building slides...")
    slide_title(prs);         print("  1/10  Title")
    slide_deliverables(prs);  print("  2/10  Deliverables status")
    slide_cost_xos(prs);      print("  3/10  Cost model — XOS")
    slide_cost_kempower(prs); print("  4/10  Cost model — Kempower")
    slide_fixed_milp(prs);    print("  5/10  Fixed DCFC MILP formulation")
    slide_kempower_milp(prs); print("  6/10  Mobile Kempower — 3 differences")
    slide_xos_simulation(prs);print("  7/10  Mobile XOS — SOC simulation")
    slide_comparison(prs);    print("  8/10  Formulation comparison table")
    slide_cost_comparison(prs);print(" 9/10  Cost comparison & gaps")
    slide_next_steps(prs);    print(" 10/10 Next steps & timeline")

    prs.save(str(OUT_PATH))
    print(f"\nSaved: {OUT_PATH}")
    print("Upload to Google Drive -> open with Google Slides (auto-converts)")


if __name__ == "__main__":
    main()
